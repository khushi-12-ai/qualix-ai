import os
import sys

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".ppt_deps"))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_THEME_COLOR

OUT = os.path.join(os.path.dirname(__file__), "Qualix_AI_Hackathon_Presentation.pptx")

# 16:9 executive dashboard language: deep navy, cyan, purple and mint.
NAVY = "071A2E"; NAVY_2 = "0D2740"; CYAN = "37D6E6"; MINT = "54E3B2"
PURPLE = "8167F5"; WHITE = "F5FAFF"; MUTED = "A9BED0"; RED = "FF6B6B"
YELLOW = "FFCE5C"; PANEL = "102E49"; GRID = "1B3D57"

def rgb(h): return RGBColor.from_string(h)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def rect(slide, x, y, w, h, color, radius=False, line=None, transparency=0):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb(color); sh.fill.transparency = transparency
    sh.line.color.rgb = rgb(line or color)
    if radius: sh.adjustments[0] = 0.12
    return sh

def text(slide, value, x, y, w, h, size=18, color=WHITE, bold=False, font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = valign
    for i, line in enumerate(value.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.alignment = align
        p.font.name = font; p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = rgb(color)
        p.space_after = Pt(2)
    tf.margin_left = tf.margin_right = Inches(.03); tf.margin_top = tf.margin_bottom = Inches(.02)
    return box

def line(slide, x1, y1, x2, y2, color=CYAN, width=1.5):
    sh = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    sh.line.color.rgb = rgb(color); sh.line.width = Pt(width); return sh

def base(title, kicker, number):
    s = prs.slides.add_slide(blank)
    rect(s, 0, 0, 13.333, 7.5, NAVY)
    # branded grid / corner detail
    for x in [9.9, 10.7, 11.5, 12.3]: line(s, x, .25, x, 7.25, GRID, .45)
    for y in [.7, 1.5, 2.3, 3.1, 3.9, 4.7, 5.5, 6.3]: line(s, 9.7, y, 13.1, y, GRID, .45)
    rect(s, .55, .52, .1, .63, CYAN, radius=True)
    text(s, kicker.upper(), .78, .5, 3, .25, 9.5, CYAN, True)
    text(s, title, .78, .76, 8.7, .56, 27, WHITE, True)
    text(s, f"QUALIX AI  /  {number:02d}", 10.05, 7.08, 2.6, .2, 8.5, MUTED, True, align=PP_ALIGN.RIGHT)
    return s

def pill(slide, label, x, y, w, color=CYAN):
    rect(slide, x, y, w, .34, color, radius=True)
    text(slide, label, x+.06, y+.075, w-.12, .16, 9, NAVY, True, align=PP_ALIGN.CENTER)

def card(slide, x, y, w, h, title, body, accent=CYAN, icon=None):
    rect(slide, x, y, w, h, PANEL, radius=True, line=GRID)
    rect(slide, x, y, .07, h, accent, radius=True)
    if icon: text(slide, icon, x+.25, y+.21, .35, .3, 16, accent, True)
    text(slide, title, x+.25+(0.42 if icon else 0), y+.21, w-.48, .27, 13, WHITE, True)
    text(slide, body, x+.25, y+.64, w-.5, h-.78, 11.3, MUTED)

# 1 cover
s = prs.slides.add_slide(blank); rect(s,0,0,13.333,7.5,NAVY)
for r, col in [(3.9,GRID),(3.0,PANEL),(2.1,"164764"),(1.2,CYAN)]:
    sh=s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.75-r/2), Inches(3.75-r/2), Inches(r), Inches(r)); sh.fill.background(); sh.line.color.rgb=rgb(col); sh.line.width=Pt(1.2)
for a,b in [(8.7,2.7),(10.5,1.6),(11.8,3.5),(10.1,5.6),(12.3,5.4)]: line(s,10.7,3.75,a,b, CYAN,.8)
for x,y in [(8.7,2.7),(10.5,1.6),(11.8,3.5),(10.1,5.6),(12.3,5.4),(10.7,3.75)]:
    sh=s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x-.07), Inches(y-.07), Inches(.14), Inches(.14)); sh.fill.solid();sh.fill.fore_color.rgb=rgb(MINT);sh.line.color.rgb=rgb(MINT)
pill(s,"HACKATHON 2026",.8,.85,1.65,MINT)
text(s,"QUALIX AI",.78,1.55,6.5,.65,38,WHITE,True)
text(s,"Secure AI Readiness\nOperating System for MSMEs",.8,2.3,6.7,1.28,27,CYAN,True)
text(s,"Transforming raw business data into a trusted foundation for AI, BI, and better decisions.",.82,3.95,5.75,.58,15,MUTED)
line(s,.82,5.0,6.5,5.0,CYAN,1)
text(s,"DATA QUALITY  •  SECURITY  •  ACTION",.82,5.24,5.8,.23,10.5,MINT,True)
text(s,"August 2026  |  MSME Data Readiness Challenge",.82,6.75,5.8,.2,10,MUTED)

# 2 Problem
s=base("The AI ambition is real. The data foundation isn’t.","Problem statement",2)
text(s,"MSME data lives across spreadsheets, POS exports, Tally ledgers and CRM dumps—often inconsistent, incomplete and unsafe to use.",.8,1.52,8.3,.55,15,MUTED)
items=[("Missing values","Revenue projections are structurally under-reported",RED),("Fuzzy duplicates","Customer counts inflate; campaigns miss the mark",YELLOW),("Schema mismatch","ERP + CRM merges fail silently",PURPLE),("No security gate","Uploaded datasets create threat exposure",CYAN)]
for i,(t,b,c) in enumerate(items): card(s,.8+(i%2)*4.35,2.35+(i//2)*1.7,3.95,1.32,t,b,c)
rect(s,9.75,1.55,2.5,4.5,PANEL,True,line=GRID)
text(s,"ROOT\nCAUSE",10.06,2.0,1.85,.7,20,CYAN,True,align=PP_ALIGN.CENTER)
line(s,10.16,2.98,11.85,2.98,MINT,1)
text(s,"No unified, automated way to diagnose raw data and translate issues into business risk.",10.05,3.28,1.95,1.25,13,WHITE,True,align=PP_ALIGN.CENTER)

# 3 solution
s=base("Qualix AI: the AI Data Doctor for MSMEs", "Proposed solution",3)
text(s,"One secure workspace to assess, improve and continuously monitor business data—without needing a data scientist.",.8,1.48,8.3,.4,15,MUTED)
features=[("01","Readiness Score","Definitive 0–100 pipeline viability score"),("02","Fix Center","Plain-language remediation for teams and executives"),("03","Smart Merge","Unified view across fragmented datasets"),("04","Trust by Design","ClamAV scan + AES-256 encryption"),("05","Executive Translation","Technical errors → business risks in 11 languages"),("06","Proof & Monitoring","Tamper-proof certificate and health alerts")]
for i,(n,t,b) in enumerate(features):
    x=.8+(i%3)*2.95;y=2.22+(i//3)*1.78
    rect(s,x,y,2.65,1.38,PANEL,True,line=GRID);text(s,n,x+.2,y+.18,.35,.25,11,CYAN,True);text(s,t,x+.2,y+.51,2.2,.23,12.5,WHITE,True);text(s,b,x+.2,y+.81,2.25,.4,10.2,MUTED)
card(s,9.82,2.22,2.55,3.16,"Outcome","From raw files to a governed, decision-ready data asset.",MINT)

# 4 scoring
s=base("One score. Five dimensions. Clear next steps.","Readiness scoring engine",4)
dimensions=[("Completeness","25%","Non-null cells",CYAN),("Consistency","25%","Formats & type conflicts",MINT),("Duplication","20%","Exact + fuzzy repeats",PURPLE),("Structure","15%","Schema adherence",YELLOW),("Anomaly","15%","IQR outliers",RED)]
for i,(n,wt,d,c) in enumerate(dimensions):
    x=.8+i*1.83
    rect(s,x,2.15,1.52,2.7,PANEL,True,line=GRID)
    text(s,wt,x+.15,2.5,1.2,.42,22,c,True,align=PP_ALIGN.CENTER)
    text(s,n,x+.13,3.2,1.25,.38,12,WHITE,True,align=PP_ALIGN.CENTER)
    text(s,d,x+.16,3.83,1.18,.42,9.5,MUTED,align=PP_ALIGN.CENTER)
    rect(s,x+.22,4.5,1.08,.07,c,True)
rect(s,.8,5.5,8.86,.82,"0B2138",True,line=GRID)
text(s,"ML Readiness = weighted base × 0.8  − 25 (target leakage)  − 15 (class imbalance >80%)",1.08,5.73,8.25,.27,13,WHITE,True,align=PP_ALIGN.CENTER)
rect(s,10.15,2.15,2.05,3.85,CYAN,True);text(s,"0–100",10.35,2.78,1.65,.55,29,NAVY,True,align=PP_ALIGN.CENTER);text(s,"Readiness\nscore",10.35,3.55,1.65,.68,18,NAVY,True,align=PP_ALIGN.CENTER);text(s,"clamped at 10–100",10.35,4.65,1.65,.25,10,NAVY,align=PP_ALIGN.CENTER)

# 5 workflow
s=base("Secure data journey: from upload to confidence", "End-to-end workflow",5)
steps=[("1","Ingest","Tally · POS · CRM · files"),("2","Secure","ClamAV scan + encryption"),("3","Unify","Smart merge + reconciliation"),("4","Profile","Stats, schema & distributions"),("5","Diagnose","5D audit + ML readiness"),("6","Act","Fixes, alerts & certificate")]
for i,(n,t,b) in enumerate(steps):
    x=.62+i*2.08
    rect(s,x,2.35,1.66,1.7,PANEL,True,line=GRID)
    rect(s,x+.61,2.04,.43,.43,CYAN,True);text(s,n,x+.61,2.125,.43,.15,11,NAVY,True,align=PP_ALIGN.CENTER)
    text(s,t,x+.14,2.77,1.38,.25,13,WHITE,True,align=PP_ALIGN.CENTER);text(s,b,x+.13,3.18,1.4,.43,9.3,MUTED,align=PP_ALIGN.CENTER)
    if i<5: line(s,x+1.68,3.19,x+2.02,3.19,MINT,2)
text(s,"Clean datasets move forward. Infected files are rejected and the user is alerted immediately.",1.15,5.22,10.85,.38,16,MINT,True,align=PP_ALIGN.CENTER)

# 6 safety trust
s=base("Zero-trust ingestion protects every dataset", "Security & governance",6)
card(s,.8,1.7,3.65,3.75,"01  Malware Defense","Every uploaded file is scanned with ClamAV before it reaches the processing layer. Infected files are rejected and logged.",RED)
card(s,4.82,1.7,3.65,3.75,"02  Encryption at Rest","Clean files are protected with Fernet AES-256 encryption. Processing occurs in memory to balance safety and speed.",CYAN)
card(s,8.84,1.7,3.65,3.75,"03  Responsible Access","PII and financial fields are classified, masked by role, and recorded in immutable audit logs.",MINT)
text(s,"Admin  →  full access     Analyst  →  PII-masked     Viewer  →  dashboard & reports",1.1,6.15,11.1,.32,13,WHITE,True,align=PP_ALIGN.CENTER)

# 7 impact translation
s=base("Technical findings become business decisions", "Multilingual impact engine",7)
text(s,"Qualix translates data-quality issues into clear executive impact statements in 11 regional languages.",.8,1.48,8.5,.35,15,MUTED)
langs=["English","Hindi","Tamil","Telugu","Kannada","Malayalam","Marathi","Bengali","Gujarati","Punjabi","Odia"]
for i,l in enumerate(langs):
    x=.82+(i%6)*1.4;y=2.15+(i//6)*.55
    rect(s,x,y,1.18,.34,"163B57",True,line=GRID);text(s,l,x+.05,y+.09,1.08,.13,8.5,WHITE,True,align=PP_ALIGN.CENTER)
rect(s,.82,3.82,7.4,1.55,PANEL,True,line=GRID)
text(s,"Example: Missing customer contacts",1.08,4.08,6.9,.25,14,CYAN,True)
text(s,"“Your customer follow-up may be incomplete, creating a risk of missed sales and lower campaign reach.”",1.08,4.51,6.75,.46,15,WHITE,False)
for i,(t,b,c) in enumerate([("Finding","What went wrong",CYAN),("Business Risk","Why it matters",RED),("Recommended Action","What to do next",MINT),("Priority","How urgently to act",YELLOW)]): card(s,8.65,1.72+i*1.02,3.45,.78,t,b,c)

# 8 tech architecture
s=base("Built as a modular, production-minded MVP", "Technical architecture",8)
layers=[("Experience","Streamlit dashboard · AI Data Doctor · role-aware views",CYAN),("API & Orchestration","FastAPI · scheduled monitoring · notification service",MINT),("Data Intelligence","Profiling · 5D quality · ML readiness · anomaly & duplicate detection",PURPLE),("Trust & Reporting","ClamAV · Fernet AES-256 · PII masking · audit log · PDF certificate",YELLOW)]
for i,(t,b,c) in enumerate(layers):
    y=1.65+i*1.12;rect(s,.82,y,8.1,.82,PANEL,True,line=GRID);rect(s,.82,y,.1,.82,c);text(s,t,1.16,y+.21,2.05,.2,13,WHITE,True);text(s,b,3.15,y+.22,5.35,.2,11,MUTED)
rect(s,9.42,1.65,2.75,4.22,"0B2138",True,line=GRID);text(s,"29",9.72,2.03,2.1,.55,35,CYAN,True,align=PP_ALIGN.CENTER);text(s,"focused service modules",9.72,2.72,2.1,.22,12,WHITE,True,align=PP_ALIGN.CENTER);line(s,9.85,3.22,11.72,3.22,MINT,1);text(s,"Designed for clear responsibility, easier testing, and future scale.",9.72,3.62,2.1,.78,12,MUTED,align=PP_ALIGN.CENTER)

# 9 demo
s=base("A simple demo that tells a complete story", "Suggested live walkthrough",9)
demos=[("Upload","Add CSV/Excel from POS, CRM or Tally","01"),("Diagnose","Inspect profile, quality dimensions and ML score","02"),("Prioritize","Open the plain-language Fix Center","03"),("Prove","Download certificate and configure monitoring","04")]
for i,(t,b,n) in enumerate(demos):
    x=.82+i*2.95;rect(s,x,2.18,2.45,2.45,PANEL,True,line=GRID);text(s,n,x+.2,2.42,.42,.2,11,CYAN,True);text(s,t,x+.2,2.95,2.0,.27,16,WHITE,True);text(s,b,x+.2,3.49,2.02,.52,11,MUTED)
    if i<3: line(s,x+2.48,3.4,x+2.85,3.4,MINT,2)
text(s,"The payoff: the team leaves with an explainable score, a fix plan, and evidence of safe handling.",1.15,5.55,10.9,.38,16,MINT,True,align=PP_ALIGN.CENTER)

# 10 challenge roadmap
s=base("Designed for today’s MSME reality—and tomorrow’s scale", "Challenges & roadmap",10)
left=[("Real-time integration","Normalize fragmented Tally, POS and CRM data."),("Schema drift","Catch changing source formats before merges break."),("Alert fatigue","Escalate meaningful changes, not every minor fluctuation.")]
right=[("Direct connectors","MySQL, PostgreSQL, MongoDB, S3 & Google Drive"),("Generative AI","Context-aware reports powered by local LLMs"),("Distributed processing","PySpark / Dask for larger-than-memory datasets")]
text(s,"NOW  /  key execution challenges",.82,1.55,4,.22,10,CYAN,True);text(s,"NEXT  /  roadmap",7.0,1.55,4,.22,10,MINT,True)
for i,(t,b) in enumerate(left): card(s,.82,2.0+i*1.23,5.35,.95,t,b,[CYAN,YELLOW,PURPLE][i])
for i,(t,b) in enumerate(right): card(s,7.0,2.0+i*1.23,5.35,.95,t,b,[MINT,CYAN,PURPLE][i])

# 11 close
s=prs.slides.add_slide(blank);rect(s,0,0,13.333,7.5,NAVY)
rect(s,.8,.95,.1,4.6,CYAN,True)
text(s,"AI adoption starts with\ndata you can trust.",1.22,1.2,7.5,1.5,34,WHITE,True)
text(s,"Qualix AI gives MSMEs a secure, explainable path from messy data to AI-ready decisions.",1.25,3.15,6.35,.6,16,MUTED)
pill(s,"SECURE  •  ACTIONABLE  •  MSME-READY",1.22,4.3,3.15,MINT)
text(s,"Thank you",1.25,5.48,3.3,.5,25,CYAN,True)
text(s,"QUALIX AI  |  Hackathon 2026",1.25,6.23,4.4,.2,10,MUTED,True)
for r,col in [(4.5,GRID),(3.5,PANEL),(2.4,"164764"),(1.4,CYAN)]:
    sh=s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5-r/2), Inches(3.72-r/2), Inches(r), Inches(r));sh.fill.background();sh.line.color.rgb=rgb(col);sh.line.width=Pt(1.2)
for x,y in [(7.6,2.3),(9.5,1.55),(11.2,2.35),(11.7,4.4),(9.4,5.85)]: line(s,9.5,3.72,x,y,CYAN,.9)

# document properties
prs.core_properties.title = "Qualix AI — Secure AI Readiness Operating System"
prs.core_properties.subject = "Hackathon Presentation"
prs.core_properties.author = "Qualix AI"
prs.save(OUT)
print(OUT)
